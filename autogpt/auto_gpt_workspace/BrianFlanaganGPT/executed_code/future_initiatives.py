import pptx

prs = pptx.Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

# Add title and subtitle
title.text = 'Future-Looking Initiatives'
subtitle.text = 'Presentation for Client'

# Save the presentation
prs.save('future_initiatives.pptx')